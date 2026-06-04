import re
import json
import asyncio
import time
from typing import Optional
import aiohttp
from bs4 import BeautifulSoup
import cohere
from .config import config
import PyPDF2
import docx
import io
import urllib.parse
from max_ai.constants import URL_PATTERN, URL_VALIDATION_PATTERN, DOMAIN_PATTERN, USER_AGENT
import logging
from .handlers import EPubHandler, CSVHandler, JSONHandler
from rich.progress import Progress, SpinnerColumn, TextColumn

logger = logging.getLogger(__name__)


try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore[misc,assignment]
try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None  # type: ignore[misc,assignment]

try:
    from mistralai import Mistral
except ImportError:
    Mistral = None  # type: ignore[misc,assignment]

try:
    from mistralai.models.chat_completion import ChatCompletionResponse
except ImportError:
    ChatCompletionResponse = None  # type: ignore[misc,assignment]


class BaseHandler:
    """
    Базовый класс для обработчиков контента.
    
    Атрибуты:
        None
    
    Методы:
        can_handle(content_type: str, url: str) -> bool:
            Проверяет, может ли обработчик обработать данный тип контента.
        
        handle(content: bytes, url: str) -> tuple[str, str]:
            Обрабатывает контент и возвращает текст и тип источника.
    """
    
    async def can_handle(self, content_type: str, url: str) -> bool:
        """
        Проверяет, может ли обработчик обработать данный тип контента.
        
        Аргументы:
            content_type (str): Тип контента.
            url (str): URL источника.
            
        Возвращает:
            bool: True, если обработчик может обработать контент, иначе False.
        """
        return False

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        """
        Обрабатывает контент и возвращает текст и тип источника.
        
        Аргументы:
            content (bytes): Содержимое файла.
            url (str): URL источника.
            
        Возвращает:
            tuple[str, str]: Текст и тип источника.
        """
        return "", ""


class PDFHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return 'application/pdf' in content_type or url.lower().endswith('.pdf')

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text[:config.get("max_content_length", 50000)], 'pdf'
        except PyPDF2.errors.PdfReadError as e:
            logger.error(f"Ошибка чтения PDF {url}: {str(e)}")
            return f"Ошибка чтения PDF {url}: {str(e)}", 'error'
        except Exception as e:
            logger.error(f"Ошибка обработки PDF {url}: {str(e)}")
            return f"Ошибка обработки PDF {url}: {str(e)}", 'error'


class DocxHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' in content_type or url.lower().endswith('.docx')

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            doc = docx.Document(io.BytesIO(content))
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text += paragraph.text + "\n"
            return text[:config.get("max_content_length", 50000)], 'docx'
        except Exception as e:
            logger.error(f"Ошибка обработки DOCX {url}: {str(e)}")
            return f"Ошибка обработки DOCX {url}: {str(e)}", 'error'


class PptxHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return ('application/vnd.openxmlformats-officedocument.presentationml.presentation' in content_type or
                url.lower().endswith('.pptx'))

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        if Presentation is None:
            logger.error(f"Ошибка: библиотека python-pptx не установлена для обработки PPTX {url}")
            return f"Ошибка: библиотека python-pptx не установлена для обработки PPTX {url}", 'error'
        try:
            prs = Presentation(io.BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text[:config.get("max_content_length", 50000)], 'pptx'
        except Exception as e:
            logger.error(f"Ошибка обработки PPTX {url}: {str(e)}")
            return f"Ошибка обработки PPTX {url}: {str(e)}", 'error'


class XlsxHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return ('application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' in content_type or
                url.lower().endswith('.xlsx'))

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        if load_workbook is None:
            logger.error(f"Ошибка: библиотека openpyxl не установлена для обработки XLSX {url}")
            return f"Ошибка: библиотека openpyxl не установлена для обработки XLSX {url}", 'error'
        try:
            wb = load_workbook(io.BytesIO(content), read_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text += row_text + "\n"
            wb.close()
            return text[:config.get("max_content_length", 50000)], 'xlsx'
        except Exception as e:
            logger.error(f"Ошибка обработки XLSX {url}: {str(e)}")
            return f"Ошибка обработки XLSX {url}: {str(e)}", 'error'


class XlsHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return ('application/vnd.ms-excel' in content_type or
                url.lower().endswith('.xls'))

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            import xlrd
            wb = xlrd.open_workbook(file_contents=content)
            text = ""
            for sheet in wb.sheets():
                for row_idx in range(sheet.nrows):
                    row = sheet.row(row_idx)
                    row_text = "\t".join(str(cell.value) for cell in row)
                    if row_text.strip():
                        text += row_text + "\n"
            return text[:config.get("max_content_length", 50000)], 'xls'
        except ImportError:
            logger.error(f"Ошибка: библиотека xlrd не установлена для обработки XLS {url}")
            return f"Ошибка: библиотека xlrd не установлена для обработки XLS {url}", 'error'
        except Exception as e:
            logger.error(f"Ошибка обработки XLS {url}: {str(e)}")
            return f"Ошибка обработки XLS {url}: {str(e)}", 'error'


class TxtHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return ('text/plain' in content_type or
                'text/markdown' in content_type or
                url.lower().endswith('.txt') or
                url.lower().endswith('.md'))

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            text = content.decode('utf-8', errors='ignore')
            return text[:config.get("max_content_length", 50000)], 'text'
        except Exception as e:
            logger.error(f"Ошибка обработки текстового файла {url}: {str(e)}")
            return f"Ошибка обработки текстового файла {url}: {str(e)}", 'error'


class YouTubeHandler(BaseHandler):
    YOUTUBE_DOMAINS = ('www.youtube.com', 'youtube.com', 'm.youtube.com', 'youtu.be')

    async def can_handle(self, content_type: str, url: str) -> bool:
        lower = url.lower()
        return any(domain in lower for domain in self.YOUTUBE_DOMAINS)

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        video_id = self._extract_video_id(url)
        if video_id:
            try:
                import yt_dlp
                with yt_dlp.YoutubeDL({'quiet': True, 'skip_download': True}) as ydl:
                    info = ydl.extract_info(url, download=False)
                    if info and isinstance(info, dict):
                        title = info.get('title')
                        description = info.get('description')
                        if title or description:
                            metadata = []
                            if title:
                                metadata.append(f"Заголовок видео: {title}")
                            if description:
                                metadata.append(f"Описание: {description}")
                            return '\n'.join(metadata)[:config.get("max_content_length", 50000)], 'youtube'
            except ImportError:
                logger.error("Библиотека yt_dlp не установлена")
                pass
            except Exception as e:
                logger.error(f"Ошибка при извлечении метаданных YouTube: {str(e)}")
                pass

            try:
                from youtube_transcript_api import YouTubeTranscriptApi
                transcript_list = None
                try:
                    transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
                except Exception as e:
                    logger.error(f"Ошибка при получении транскрипта YouTube: {str(e)}")
                    try:
                        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
                        transcript_list = transcript_list.find_transcript(['ru', 'en']).fetch()
                    except Exception as e:
                        logger.error(f"Ошибка при поиске транскрипта YouTube: {str(e)}")
                        transcript_list = None

                if transcript_list:
                    text = ' '.join([item['text'] for item in transcript_list])
                    return text[:config.get("max_content_length", 50000)], 'youtube'
            except ImportError:
                logger.error("Библиотека youtube_transcript_api не установлена")
                pass
            except Exception as e:
                logger.error(f"Ошибка при обработке транскрипта YouTube: {str(e)}")
                pass

        try:
            html = content.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html, 'html.parser')
            title = None
            description = None

            for selector in [
                ("meta", {"property": "og:title"}),
                ("meta", {"name": "twitter:title"}),
            ]:
                tag = soup.find(*selector)
                if tag and tag.get('content'):
                    title = tag['content'].strip()
                    break

            for selector in [
                ("meta", {"property": "og:description"}),
                ("meta", {"name": "twitter:description"}),
                ("meta", {"name": "description"}),
            ]:
                tag = soup.find(*selector)
                if tag and tag.get('content'):
                    description = tag['content'].strip()
                    break

            json_data = self._extract_yt_initial_data(html)
            if json_data is not None:
                video_details = json_data.get('videoDetails', {})
                title = title or video_details.get('title')
                description = description or video_details.get('shortDescription')

            items = []
            if title:
                items.append(f"og:title: {title}")
            if description:
                items.append(f"shortDescription: {description}")

            if not items:
                return f"Не удалось извлечь og:title или shortDescription для {url}", 'youtube'

            result = '\n'.join(items)[:config.get("max_content_length", 50000)]
            return result, 'youtube'
        except Exception as e:
            return f"Ошибка обработки YouTube {url}: {str(e)}", 'error'

    @staticmethod
    def _extract_video_id(url: str) -> Optional[str]:
        parsed = urllib.parse.urlparse(url)
        if parsed.netloc in ('www.youtube.com', 'youtube.com', 'm.youtube.com'):
            query = urllib.parse.parse_qs(parsed.query)
            return query.get('v', [None])[0]
        if parsed.netloc == 'youtu.be':
            return parsed.path.lstrip('/')
        return None

    @staticmethod
    def _extract_yt_initial_data(html: str) -> Optional[dict]:
        patterns = [
            r"ytInitialPlayerResponse\s*=\s*({.+?})\s*;",
            r"ytInitialPlayerResponse\s*:\s*({.+?})\s*,\s*" ,
            r'"ytInitialPlayerResponse"\s*:\s*({.+?})',
        ]
        for pattern in patterns:
            match = re.search(pattern, html, flags=re.DOTALL)
            if match:
                try:
                    text = match.group(1)
                    return json.loads(text)
                except Exception as e:
                    logger.error(f"Ошибка при парсинге JSON из HTML: {str(e)}")
                    continue
        return None


class HTMLHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return True

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            html = content.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html, 'html.parser')
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text(separator='\n', strip=True)
            return text[:config.get("max_content_length", 50000)], 'web'
        except UnicodeDecodeError as e:
            logger.error(f"Ошибка декодирования страницы {url}: {str(e)}")
            return f"Ошибка декодирования страницы {url}: {str(e)}", 'error'
        except Exception as e:
            logger.error(f"Ошибка обработки HTML {url}: {str(e)}")
            return f"Ошибка обработки HTML {url}: {str(e)}", 'error'


class AIAgent:
    """
    Основной класс AI-агента для обработки запросов и источников.
    
    Атрибуты:
        cohere_client: Клиент для работы с API Cohere.
        mistral_client: Клиент для работы с API Mistral.
        model_name (str): Название модели Cohere.
        mistral_model (str): Название модели Mistral.
        use_mistral (bool): Флаг использования Mistral.
        handlers (list[BaseHandler]): Список обработчиков для различных типов контента.
        conversation_history (list[dict[str, str]]): История диалога.
        _rate_limiter: Семафор для ограничения скорости запросов.
        _timeout (int): Таймаут для запросов.
        _max_retries (int): Максимальное количество попыток для запросов.
        _url_cache (dict): Кэш для хранения результатов обработки URL.
    
    Методы:
        clear_history(): Очищает историю диалога.
        clear_cache(): Очищает кэш URL.
        _summarize_content(content: str, source_type: str) -> str:
            Суммаризирует контент, если он слишком длинный.
        extract_urls(text: str) -> list[str]:
            Извлекает URL из текста.
        _fetch_url(url: str, session: aiohttp.ClientSession, retries: int) -> tuple[str, str]:
            Загружает и обрабатывает контент по URL.
        fetch_urls_async(urls: list[str]) -> list[tuple[str, str]]:
            Асинхронно загружает и обрабатывает контент по списку URL.
        _run_async(coro) -> list[tuple[str, str]]:
            Запускает асинхронную корутину.
        run(query: str, sources: Optional[list[str]]) -> tuple[str, int]:
            Выполняет запрос и возвращает ответ и количество использованных токенов.
    """
    
    def __init__(self, cohere_key: Optional[str] = None, mistral_key: Optional[str] = None, model: Optional[str] = None, use_mistral: bool = True) -> None:
        cohere_key_str = config.get_cohere_key(cohere_key)
        if not cohere_key_str:
            raise ValueError("Cohere API key is required. Set COHERE_API_KEY or pass --cohere-key.")
        self.cohere_client = cohere.ClientV2(cohere_key_str)
        self.model_name = model or config.get("cohere_model", "command-a-03-2025")
        self.mistral_model = config.get("mistral_model", "mistral-large-latest")
        self.use_mistral = use_mistral
        self.mistral_client: Optional[Mistral] = None
        mistral_api_key = config.get_mistral_key(mistral_key) if self.use_mistral else None
        if self.use_mistral and mistral_api_key and Mistral is not None:
            try:
                self.mistral_client = Mistral(api_key=mistral_api_key)
            except Exception as e:
                logger.error(f"Ошибка при инициализации Mistral клиента: {str(e)}")
                self.mistral_client = None
        self.handlers: list[BaseHandler] = [
            PDFHandler(),
            DocxHandler(),
            PptxHandler(),
            XlsxHandler(),
            XlsHandler(),
            TxtHandler(),
            YouTubeHandler(),
            HTMLHandler(),
            EPubHandler(),
            CSVHandler(),
            JSONHandler()
        ]
        self.conversation_history: list[dict[str, str]] = []
        self._rate_limiter = asyncio.Semaphore(config.get("rate_limit", 5))
        self._timeout = config.get("timeout", 30)
        self._max_retries = config.get("max_retries", 3)
        self._url_cache = {}

    def clear_history(self) -> None:
        """
        Очищает историю диалога.
        """
        self.conversation_history = []

    def clear_cache(self) -> None:
        """
        Очищает кэш URL.
        """
        self._url_cache.clear()
        logger.info("Кэш URL очищен")

    def _summarize_content(self, content: str, source_type: str) -> str:
        """
        Суммаризирует контент, если он слишком длинный.
        
        Аргументы:
            content (str): Текст для суммаризации.
            source_type (str): Тип источника.
            
        Возвращает:
            str: Суммаризированный текст или обрезанный текст, если суммаризация не удалась.
        """
        threshold = config.get("summarize_threshold", 40000)
        if len(content) <= threshold:
            return content
        
        text_to_summarize = content[:threshold]
        summary_query = f"Суммаризируй следующий текст (тип: {source_type}):\n\n{text_to_summarize}"
        
        try:
            summary_response = self.cohere_client.chat(
                model=self.model_name,
                messages=[{"role": "user", "content": summary_query}]
            )
            if summary_response.message.content and summary_response.message.content[0].text:
                logger.info(f"Контент успешно суммаризирован (тип: {source_type})")
                return summary_response.message.content[0].text + "\n\n[Текст был автоматически суммаризирован]"
        except Exception as e:
            logger.error(f"Ошибка при суммаризации контента: {str(e)}")
        
        logger.warning(f"Суммаризация не удалась, возвращаем обрезанный контент (тип: {source_type})")
        return content[:config.get("max_content_length", 50000)]

    def extract_urls(self, text: str) -> list[str]:
        """
        Извлекает URL из текста.
        
        Аргументы:
            text (str): Текст, из которого нужно извлечь URL.
            
        Возвращает:
            list[str]: Список извлеченных URL.
        """
        urls = URL_PATTERN.findall(text)
        return [url for url in urls if URL_VALIDATION_PATTERN.match(url)]

    async def _fetch_url(self, url: str, session: aiohttp.ClientSession, retries: int = 3) -> tuple[str, str]:
        """
        Загружает и обрабатывает контент по URL.
        
        Аргументы:
            url (str): URL для загрузки.
            session (aiohttp.ClientSession): Сессия для выполнения запросов.
            retries (int): Количество попыток загрузки.
            
        Возвращает:
            tuple[str, str]: Текст и тип источника или сообщение об ошибке.
        """
        if not URL_VALIDATION_PATTERN.match(url):
            logger.error(f"Невалидный URL: {url}")
            return f"Невалидный URL: {url}", 'error'
        
        if url in self._url_cache:
            logger.info(f"Используем кэшированный результат для URL: {url}")
            return self._url_cache[url]
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept-Language': 'en-US,en;q=0.9,ru;q=0.8',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8'
        }

        last_error: Optional[Exception] = None
        for attempt in range(self._max_retries):
            try:
                async with self._rate_limiter:
                    async with session.get(url, timeout=aiohttp.ClientTimeout(total=self._timeout), headers=headers) as response:
                        response.raise_for_status()
                        content = await response.read()
                        content_type = response.headers.get('content-type', '').lower()

                        for handler in self.handlers:
                            if await handler.can_handle(content_type, url):
                                result = await handler.handle(content, url)
                                self._url_cache[url] = result
                                return result

                        logger.warning(f"Не удалось обработать содержимое {url}: неизвестный тип контента")
                        return f"Не удалось обработать содержимое {url}: неизвестный тип контента", 'error'
            except aiohttp.ClientError as e:
                last_error = e
                logger.error(f"Ошибка при загрузке URL {url} (попытка {attempt + 1}): {str(e)}")
            except Exception as e:
                last_error = e
                logger.error(f"Ошибка при обработке URL {url} (попытка {attempt + 1}): {str(e)}")
            
            if attempt < self._max_retries - 1:
                await asyncio.sleep(0.5 * (2 ** attempt))

        logger.error(f"Ошибка загрузки {url} после {self._max_retries} попыток: {str(last_error)}")
        return f"Ошибка загрузки {url} после {self._max_retries} попыток: {str(last_error)}", 'error'

    async def fetch_urls_async(self, urls: list[str]) -> list[tuple[str, str]]:
        """
        Асинхронно загружает и обрабатывает контент по списку URL.
        
        Аргументы:
            urls (list[str]): Список URL для загрузки.
            
        Возвращает:
            list[tuple[str, str]]: Список результатов обработки URL.
        """
        async with aiohttp.ClientSession() as session:
            tasks = [self._fetch_url(url, session) for url in urls]
            return await asyncio.gather(*tasks)

    def _run_async(self, coro) -> list[tuple[str, str]]:
        """
        Запускает асинхронную корутину.
        
        Аргументы:
            coro: Корутина для выполнения.
            
        Возвращает:
            list[tuple[str, str]]: Результат выполнения корутины.
        """
        try:
            return asyncio.run(coro)
        except RuntimeError:
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                return pool.submit(asyncio.run, coro).result()

    def run(self, query: str, sources: Optional[list[str]] = None) -> tuple[str, int]:
        """
        Выполняет запрос и возвращает ответ и количество использованных токенов.
        
        Аргументы:
            query (str): Запрос пользователя.
            sources (Optional[list[str]]): Список источников для обработки.
            
        Возвращает:
            tuple[str, int]: Ответ и количество использованных токенов.
        """
        sources = sources or []
        urls = list(dict.fromkeys(self.extract_urls(query) + [url for url in sources if URL_VALIDATION_PATTERN.match(url)]))
        content_parts: list[str] = []
        if urls:
            with Progress(
                SpinnerColumn(),
                TextColumn("[progress.description]{task.description}"),
                transient=True,
            ) as progress:
                progress.add_task(description="Обработка источников...", total=None)
                results = self._run_async(self.fetch_urls_async(urls))
                for url, (content, source_type) in zip(urls, results):
                    if source_type != 'error':
                        content = self._summarize_content(content, source_type)
                    content_parts.append(f"\n\n--- Содержание {url} (тип: {source_type}) ---\n{content}")

        if urls:
            full_query = (
                "Прочитай следующие источники (если это видео — используй транскрипт) и ответь на вопрос:\n"
                + "\n".join(content_parts)
                + f"\n\nВопрос: {query}"
            )
        else:
            full_query = query

        history_parts = []
        for msg in self.conversation_history[-10:]:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            history_parts.append(f"{role}: {content}")
        history_text = "\n".join(history_parts)
        if history_text:
            full_query = f"История диалога:\n{history_text}\n\nТекущий запрос:\n{full_query}"

        cohere_tokens = 0
        draft = None
        if self.mistral_client is None:
            try:
                cohere_response = self.cohere_client.chat(
                    model=self.model_name,
                    messages=[{"role": "user", "content": full_query}]
                )
                draft = cohere_response.message.content[0].text
                if hasattr(cohere_response, 'usage') and cohere_response.usage:
                    tokens_obj = getattr(cohere_response.usage, 'tokens', None)
                    if tokens_obj:
                        cohere_tokens = int(tokens_obj.input_tokens or 0) + int(tokens_obj.output_tokens or 0)
            except cohere.errors.CohereError as e:
                logger.error(f"Ошибка при вызове Cohere: {str(e)}")
                return f"Ошибка при вызове Cohere: {str(e)}", 0
            except Exception as e:
                logger.error(f"Ошибка при вызове Cohere: {str(e)}")
                return f"Ошибка при вызове Cohere: {str(e)}", 0

        mistral_tokens = 0
        response = None
        if self.mistral_client is not None:
            mistral_prompt = (
                "Ты — опытный ассистент, читающий ссылки, транскрипты видео и тексты страниц, и дающий точный ответ по запросу пользователя."
                "\n\nИспользуй только информацию, полученную из предоставленного контента."
                f"\n\n{full_query}\n\n"
                "Ответь коротко, ясно и по существу."
            )
            try:
                mistral_response = self.mistral_client.chat(
                    model=self.mistral_model,
                    messages=[{"role": "user", "content": mistral_prompt}]
                )
                if hasattr(mistral_response, 'usage') and mistral_response.usage:
                    mistral_tokens = mistral_response.usage.total_tokens or 0
                if hasattr(mistral_response, 'choices') and mistral_response.choices:
                    response = mistral_response.choices[0].message.content
                elif hasattr(mistral_response, 'message') and hasattr(mistral_response.message, 'content'):
                    improved = mistral_response.message.content
                    if isinstance(improved, list) and improved:
                        response = improved[0].text if hasattr(improved[0], 'text') else str(improved)
                    else:
                        response = str(improved)
                else:
                    response = str(mistral_response)
            except Exception as e:
                logger.error(f"Ошибка при вызове Mistral: {str(e)}")
                response = None

        if response is None:
            if draft is None:
                try:
                    cohere_response = self.cohere_client.chat(
                        model=self.model_name,
                        messages=[{"role": "user", "content": full_query}]
                    )
                    draft = cohere_response.message.content[0].text
                    if hasattr(cohere_response, "usage") and cohere_response.usage:
                        tokens_obj = getattr(cohere_response.usage, "tokens", None)
                        if tokens_obj:
                            cohere_tokens = int(tokens_obj.input_tokens or 0) + int(tokens_obj.output_tokens or 0)
                except cohere.errors.CohereError as e:
                    logger.error(f"Ошибка при вызове Cohere: {str(e)}")
                    return f"Ошибка при вызове Cohere: {str(e)}", 0
                except Exception as e:
                    logger.error(f"Ошибка при вызове Cohere: {str(e)}")
                    return f"Ошибка при вызове Cohere: {str(e)}", 0
            response = draft

        self.conversation_history.append({"role": "user", "content": query})
        self.conversation_history.append({"role": "assistant", "content": response})
        return response, cohere_tokens + mistral_tokens