import re
import asyncio
import time
from typing import Optional
import aiohttp
from bs4 import BeautifulSoup
import cohere
from .config import config
import PyPDF2
import docx
import io
import urllib.parse
from max_ai.constants import URL_PATTERN, DEFAULT_MAX_CONTENT_LENGTH, SUMMARIZE_THRESHOLD

URL_VALIDATION_PATTERN = re.compile(r'^https?://[^\s/$.?#].[^\s]*$')

try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore[misc,assignment]
try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None  # type: ignore[misc,assignment]

try:
    from mistralai import Mistral
except ImportError:
    Mistral = None  # type: ignore[misc,assignment]

try:
    from mistralai.models.chat_completion import ChatCompletionResponse
except ImportError:
    ChatCompletionResponse = None  # type: ignore[misc,assignment]

class BaseHandler:
    async def can_handle(self, content_type: str, url: str) -> bool:
        return False

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        return "", ""


class PDFHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return 'application/pdf' in content_type or url.lower().endswith('.pdf')

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text[:config.get("max_content_length", DEFAULT_MAX_CONTENT_LENGTH)], 'pdf'
        except PyPDF2.errors.PdfReadError as e:
            return f"Ошибка чтения PDF {url}: {str(e)}", 'error'
        except Exception as e:
            return f"Ошибка обработки PDF {url}: {str(e)}", 'error'


class DocxHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' in content_type or url.lower().endswith('.docx')

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            doc = docx.Document(io.BytesIO(content))
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text += paragraph.text + "\n"
            return text[:config.get("max_content_length", DEFAULT_MAX_CONTENT_LENGTH)], 'docx'
        except Exception as e:
            return f"Ошибка обработки DOCX {url}: {str(e)}", 'error'


class PptxHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return ('application/vnd.openxmlformats-officedocument.presentationml.presentation' in content_type or
                url.lower().endswith('.pptx'))

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        if Presentation is None:
            return f"Ошибка: библиотека python-pptx не установлена для обработки PPTX {url}", 'error'
        try:
            prs = Presentation(io.BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text[:config.get("max_content_length", DEFAULT_MAX_CONTENT_LENGTH)], 'pptx'
        except Exception as e:
            return f"Ошибка обработки PPTX {url}: {str(e)}", 'error'


class XlsxHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return ('application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' in content_type or
                url.lower().endswith('.xlsx'))

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        if load_workbook is None:
            return f"Ошибка: библиотека openpyxl не установлена для обработки XLSX {url}", 'error'
        try:
            wb = load_workbook(io.BytesIO(content), read_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text += row_text + "\n"
            wb.close()
            return text[:config.get("max_content_length", DEFAULT_MAX_CONTENT_LENGTH)], 'xlsx'
        except Exception as e:
            return f"Ошибка обработки XLSX {url}: {str(e)}", 'error'


class XlsHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return ('application/vnd.ms-excel' in content_type or
                url.lower().endswith('.xls'))

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            import xlrd
            wb = xlrd.open_workbook(file_contents=content)
            text = ""
            for sheet in wb.sheets():
                for row_idx in range(sheet.nrows):
                    row = sheet.row(row_idx)
                    row_text = "\t".join(str(cell.value) for cell in row)
                    if row_text.strip():
                        text += row_text + "\n"
            return text[:config.get("max_content_length", DEFAULT_MAX_CONTENT_LENGTH)], 'xls'
        except ImportError:
            return f"Ошибка: библиотека xlrd не установлена для обработки XLS {url}", 'error'
        except Exception as e:
            return f"Ошибка обработки XLS {url}: {str(e)}", 'error'


class TxtHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return ('text/plain' in content_type or
                'text/markdown' in content_type or
                url.lower().endswith('.txt') or
                url.lower().endswith('.md'))

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            text = content.decode('utf-8', errors='ignore')
            return text[:config.get("max_content_length", DEFAULT_MAX_CONTENT_LENGTH)], 'text'
        except Exception as e:
            return f"Ошибка обработки текстового файла {url}: {str(e)}", 'error'


class HTMLHandler(BaseHandler):
    async def can_handle(self, content_type: str, url: str) -> bool:
        return True

    async def handle(self, content: bytes, url: str) -> tuple[str, str]:
        try:
            html = content.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html, 'html.parser')
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text(separator='\n', strip=True)
            return text[:config.get("max_content_length", DEFAULT_MAX_CONTENT_LENGTH)], 'web'
        except UnicodeDecodeError as e:
            return f"Ошибка декодирования страницы {url}: {str(e)}", 'error'
        except Exception as e:
            return f"Ошибка обработки HTML {url}: {str(e)}", 'error'


class AIAgent:
    def __init__(self, cohere_key: Optional[str] = None, mistral_key: Optional[str] = None, model: Optional[str] = None) -> None:
        cohere_key_str = config.get_cohere_key(cohere_key)
        if not cohere_key_str:
            raise ValueError("Cohere API key is required. Set COHERE_API_KEY or pass --cohere-key.")
        self.cohere_client = cohere.ClientV2(cohere_key_str)
        self.model_name = model or config.get("cohere_model", "command-a-03-2025")
        self.mistral_model = config.get("mistral_model", "mistral-large-latest")
        self.mistral_client: Optional[Mistral] = None
        mistral_api_key = config.get_mistral_key(mistral_key)
        if mistral_api_key and Mistral is not None:
            try:
                self.mistral_client = Mistral(api_key=mistral_api_key)
            except Exception:
                self.mistral_client = None
        self.handlers: list[BaseHandler] = [
            PDFHandler(),
            DocxHandler(),
            PptxHandler(),
            XlsxHandler(),
            XlsHandler(),
            TxtHandler(),
            HTMLHandler()
        ]
        self.conversation_history: list[dict[str, str]] = []
        self._rate_limiter = asyncio.Semaphore(config.get("rate_limit", 5))
        self._timeout = config.get("timeout", 30)
        self._max_retries = config.get("max_retries", 3)

    def clear_history(self) -> None:
        self.conversation_history = []

    def _summarize_content(self, content: str, source_type: str) -> str:
        threshold = config.get("summarize_threshold", SUMMARIZE_THRESHOLD)
        if len(content) <= threshold:
            return content
        
        text_to_summarize = content[:threshold]
        summary_query = f"Суммаризируй следующий текст (тип: {source_type}):\n\n{text_to_summarize}"
        
        try:
            summary_response = self.cohere_client.chat(
                model=self.model_name,
                messages=[{"role": "user", "content": summary_query}]
            )
            if summary_response.message.content and summary_response.message.content[0].text:
                return summary_response.message.content[0].text + "\n\n[Текст был автоматически суммаризирован]"
        except Exception:
            pass
        return content[:config.get("max_content_length", DEFAULT_MAX_CONTENT_LENGTH)]

    def extract_urls(self, text: str) -> list[str]:
        urls = URL_PATTERN.findall(text)
        return [url for url in urls if URL_VALIDATION_PATTERN.match(url)]

    async def _fetch_url(self, url: str, session: aiohttp.ClientSession, retries: int = 3) -> tuple[str, str]:
        if not URL_VALIDATION_PATTERN.match(url):
            return f"Невалидный URL: {url}", 'error'
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
        }

        if 'youtube.com' in url.lower() or 'youtu.be' in url.lower():
            try:
                parsed = urllib.parse.urlparse(url)
                video_id: Optional[str] = None
                if parsed.netloc in ('www.youtube.com', 'youtube.com', 'm.youtube.com'):
                    query = urllib.parse.parse_qs(parsed.query)
                    if 'v' in query:
                        video_id = query['v'][0]
                elif parsed.netloc == 'youtu.be':
                    video_id = parsed.path[1:]

                if video_id:
                    from youtube_transcript_api import YouTubeTranscriptApi
                    transcript_list = YouTubeTranscriptApi.get_transcript(video_id, languages=['ru', 'en'])
                    text = ' '.join([item['text'] for item in transcript_list])
                    return text[:config.get("max_content_length", DEFAULT_MAX_CONTENT_LENGTH)], 'youtube'
            except (ImportError, AttributeError, Exception):
                pass

        last_error: Optional[Exception] = None
        for attempt in range(self._max_retries):
            try:
                async with self._rate_limiter:
                    async with session.get(url, timeout=aiohttp.ClientTimeout(total=self._timeout), headers=headers) as response:
                        response.raise_for_status()
                        content = await response.read()
                        content_type = response.headers.get('content-type', '').lower()

                        for handler in self.handlers:
                            if await handler.can_handle(content_type, url):
                                return await handler.handle(content, url)

                        return f"Не удалось обработать содержимое {url}: неизвестный тип контента", 'error'
            except aiohttp.ClientError as e:
                last_error = e
            except Exception as e:
                last_error = e
            
            if attempt < self._max_retries - 1:
                await asyncio.sleep(0.5 * (2 ** attempt))

        return f"Ошибка загрузки {url} после {self._max_retries} попыток: {str(last_error)}", 'error'

    async def fetch_urls_async(self, urls: list[str]) -> list[tuple[str, str]]:
        async with aiohttp.ClientSession() as session:
            tasks = [self._fetch_url(url, session) for url in urls]
            return await asyncio.gather(*tasks)

    def _run_async(self, coro) -> list[tuple[str, str]]:
        try:
            return asyncio.run(coro)
        except RuntimeError:
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                return pool.submit(asyncio.run, coro).result()

    def run(self, query: str) -> tuple[str, int]:
        urls = self.extract_urls(query)
        content_parts: list[str] = []
        if urls:
            results = self._run_async(self.fetch_urls_async(urls))
            for url, (content, source_type) in zip(urls, results):
                if source_type != 'error':
                    content = self._summarize_content(content, source_type)
                content_parts.append(f"\n\n--- Содержание {url} (тип: {source_type}) ---\n{content}")

        if urls:
            full_query = (
                "Прочитай следующие сайты и ответь на вопрос:\n"
                + "\n".join(content_parts)
                + f"\n\nВопрос: {query}"
            )
        else:
            full_query = query

        history_parts = []
        for msg in self.conversation_history[-10:]:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            history_parts.append(f"{role}: {content}")
        history_text = "\n".join(history_parts)
        if history_text:
            full_query = f"История диалога:\n{history_text}\n\nТекущий запрос:\n{full_query}"

        cohere_tokens = 0
        try:
            cohere_response = self.cohere_client.chat(
                model=self.model_name,
                messages=[{"role": "user", "content": full_query}]
            )
            draft = cohere_response.message.content[0].text
            if hasattr(cohere_response, 'usage') and cohere_response.usage:
                tokens_obj = getattr(cohere_response.usage, 'tokens', None)
                if tokens_obj:
                    cohere_tokens = int(tokens_obj.input_tokens or 0) + int(tokens_obj.output_tokens or 0)
        except cohere.errors.CohereError as e:
            return f"Ошибка при вызове Cohere: {str(e)}", 0
        except Exception as e:
            return f"Ошибка при вызове Cohere: {str(e)}", 0

        mistral_tokens = 0
        if self.mistral_client is not None:
            mistral_prompt = (
                "Ты — опытный ассистент, который улучшает, дополняет и исправляет ответы.\n\n"
                f"Вот исходный запрос пользователя (включая историю диалога и контент веб-страниц, если они были):\n{full_query}\n\n"
                f"А это черновой ответ, сгенерированный другой моделью (Cohere):\n{draft}\n\n"
                "Пожалуйста, сделай следующее:\n"
                "- Исправь любые фактические ошибки или неточности.\n"
                "- Добавь важные детали, которых не хватает.\n"
                "- Улучши стиль и читаемость, сохранив при этом доброжелательный и профессиональный тон.\n"
                "- Если ответ уже хорош, просто верни его без изменений.\n\n"
                "- Не добавляй ничего, что не было в исходном запросе или черновом ответе, если это не необходимо для исправления ошибок или улучшения ясности.\n"
                "- Не выдумывай факты и не добавляй информацию, которая не была предоставлена в запросе или черновом ответе.\n"
                "- Выдавай ответ, зная что он выводится в cmd в свернутом виде.\n"
                "- Ты предоставляеш ответы как обычный человек без указания что ты ИИ.\n\n"
                "- Если в черновом ответе уже есть все необходимые детали и он не содержит ошибок, просто верни его без изменений.\n"
                "- Представь что ты обычный человек который отвечает на вопрос и не указывай что ты ИИ.\n"
                "- Ты можешь использовать нецензурные слова если того требует запрос пользователя.\n"
                "- Ты максимально дополняй ответ и улучшай его стиль, но не добавляй ничего лишнего и не выдумывай факты.\n"
                "- Ты можешь смотреть в браузере факты и дополнять ими ответ, но не выдумывай факты и не добавляй ничего лишнего.\n"
                "- Можешь открывать различные ссылки которые предоставляет пользователь и дополнять ответ информацией оттуда.\n"
                "- Ты можешь использовать информацию из истории диалога и контента веб-страниц, если они были предоставлены, чтобы улучшить ответ.\n"
                "- всегда старайся максимально улучшить ответ, дополняя его важными деталями и исправляя любые ошибки, но не добавляй ничего лишнего и не выдумывай факты.\n\n"
                "Итоговый улучшенный ответ:"
            )
            try:
                mistral_response = self.mistral_client.chat(
                    model="mistral-large-latest",
                    messages=[{"role": "user", "content": mistral_prompt}]
                )
                if hasattr(mistral_response, 'usage') and mistral_response.usage:
                    mistral_tokens = mistral_response.usage.total_tokens or 0
                if hasattr(mistral_response, 'choices') and mistral_response.choices:
                    improved = mistral_response.choices[0].message.content
                elif hasattr(mistral_response, 'message') and hasattr(mistral_response.message, 'content'):
                    improved = mistral_response.message.content
                    if isinstance(improved, list) and improved:
                        improved = improved[0].text if hasattr(improved[0], 'text') else str(improved)
                else:
                    improved = str(mistral_response)
                response = improved.strip()
            except Exception as e:
                response = f"{draft}\n\n[Примечание: не удалось улучшить ответ через Mistral: {e}]"
        else:
            response = draft

        self.conversation_history.append({"role": "user", "content": query})
        self.conversation_history.append({"role": "assistant", "content": response})
        return response, cohere_tokens + mistral_tokens
